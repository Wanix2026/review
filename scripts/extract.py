#!/usr/bin/env python3
"""统一文件提取脚本：支持 PDF、DOCX、DOC、PPTX、XLSX"""

import sys, os, subprocess, json, hashlib, shutil
from pathlib import Path

def sha256(filepath):
    h = hashlib.sha256()
    with open(filepath, 'rb') as f:
        while chunk := f.read(8192):
            h.update(chunk)
    return h.hexdigest()

def diagnose_pdf(filepath):
    """诊断PDF类型"""
    try:
        result = subprocess.run(
            ['pdf2txt', '-p', '1', filepath],
            capture_output=True, text=True, timeout=30
        )
        text = result.stdout.strip()
        if result.returncode != 0 and 'encrypted' in result.stderr.lower():
            return 'encrypted', ''
        char_count = len(text.replace('\n', '').replace(' ', ''))
        if char_count < 50:
            return 'scan', text
        return 'text', text
    except Exception as e:
        return 'error', str(e)

def extract_pdf(filepath, output_path):
    """提取PDF文本"""
    ptype, sample = diagnose_pdf(filepath)

    if ptype == 'encrypted':
        return {'status': 'encrypted', 'error': 'PDF已加密，需手动解密后重新放入'}

    if ptype == 'scan':
        # 尝试OCR
        try:
            result = subprocess.run(
                ['pdf2txt', '-o', output_path, filepath],
                capture_output=True, text=True, timeout=120
            )
            text = Path(output_path).read_text(encoding='utf-8', errors='replace')
            if len(text.strip()) < 100:
                return {'status': 'scan_low_quality',
                        'error': 'OCR提取文字过少，可能为纯扫描图片PDF'}
            return {'status': 'ok_ocr', 'pages': text.count('\f') + 1}
        except Exception as e:
            return {'status': 'error', 'error': str(e)}

    try:
        result = subprocess.run(
            ['pdf2txt', '-o', output_path, filepath],
            capture_output=True, text=True, timeout=120
        )
        text = Path(output_path).read_text(encoding='utf-8', errors='replace')
        pages = text.count('\f') + 1
        chars = len(text.replace('\n','').replace(' ',''))

        # 乱码检测
        garbled = sum(1 for c in text if ord(c) > 0x4e00 and ord(c) < 0x9fff)
        total_chinese = sum(1 for c in text if '一' <= c <= '鿿' or '㐀' <= c <= '䶿')
        garbled_ratio = (len(text) - chars) / max(len(text), 1) if chars < len(text) * 0.3 else 0

        return {
            'status': 'ok',
            'pages': pages,
            'chars': chars,
            'garbled_ratio': round(garbled_ratio, 3)
        }
    except Exception as e:
        return {'status': 'error', 'error': str(e)}

def extract_pptx(filepath, output_path):
    """提取PPTX文本"""
    try:
        code = f'''
import sys
sys.path.insert(0, '{shutil.which("python3")}')
from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation("{filepath}")
lines = []
slide_warnings = []

for i, slide in enumerate(prs.slides, 1):
    lines.append(f"\\n--- Slide {{i}} ---")
    has_text = False
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    lines.append(t)
                    has_text = True
        if shape.has_table:
            table = shape.table
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                lines.append(" | ".join(cells))
                has_text = True
        if hasattr(shape, 'image') and shape.image:
            slide_warnings.append(f"⚠️ Slide {{i}}: [图片]")
    if not has_text:
        lines.append("[此页无文字内容]")
        slide_warnings.append(f"⚠️ Slide {{i}}: [无文字内容/可能为SmartArt]")

with open("{output_path}", 'w') as f:
    f.write('\\n'.join(lines))
    if slide_warnings:
        f.write('\\n\\n## 提取警告\\n')
        f.write('\\n'.join(slide_warnings))
'''
        result = subprocess.run(
            ['python3', '-c', code],
            capture_output=True, text=True, timeout=60
        )
        if result.returncode != 0:
            if 'No module named' in result.stderr:
                return {'status': 'error', 'error': 'python-pptx 未安装，请 pip3 install python-pptx'}
            return {'status': 'error', 'error': result.stderr[:500]}

        text = Path(output_path).read_text(encoding='utf-8', errors='replace')
        slides = text.count('--- Slide')
        return {
            'status': 'ok',
            'slides': slides,
            'chars': len(text),
            'warnings': [l for l in text.split('\n') if l.startswith('⚠️')]
        }
    except Exception as e:
        return {'status': 'error', 'error': str(e)}

def extract_docx(filepath, output_path):
    """提取DOCX文本"""
    try:
        code = f'''
from docx import Document
doc = Document("{filepath}")
lines = []
for para in doc.paragraphs:
    if para.text.strip():
        lines.append(para.text)
for table in doc.tables:
    for row in table.rows:
        cells = [cell.text.strip() for cell in row.cells]
        lines.append(" | ".join(cells))
with open("{output_path}", 'w') as f:
    f.write('\\n'.join(lines))
'''
        result = subprocess.run(
            ['python3', '-c', code],
            capture_output=True, text=True, timeout=30
        )
        if result.returncode != 0:
            return {'status': 'error', 'error': result.stderr[:500]}
        text = Path(output_path).read_text(encoding='utf-8', errors='replace')
        return {'status': 'ok', 'chars': len(text)}
    except Exception as e:
        return {'status': 'error', 'error': str(e)}

def extract_doc(filepath, output_path):
    """提取DOC文本（旧格式）"""
    try:
        result = subprocess.run(
            ['textutil', '-convert', 'txt', '-output', output_path, filepath],
            capture_output=True, text=True, timeout=30
        )
        if result.returncode != 0:
            return {'status': 'error', 'error': 'textutil转换失败，请将.doc另存为.docx'}
        text = Path(output_path).read_text(encoding='utf-8', errors='replace')
        return {'status': 'ok', 'chars': len(text)}
    except Exception as e:
        return {'status': 'error', 'error': str(e)}

def extract_xlsx(filepath, output_path):
    """提取XLSX文本"""
    try:
        code = f'''
from openpyxl import load_workbook
wb = load_workbook("{filepath}", data_only=True)
lines = []
for name in wb.sheetnames:
    ws = wb[name]
    lines.append(f"\\n--- Sheet: {{name}} ---")
    for row in ws.iter_rows(values_only=True):
        cells = [str(c) if c is not None else "" for c in row]
        if any(c.strip() for c in cells):
            lines.append(" | ".join(cells))
with open("{output_path}", 'w') as f:
    f.write('\\n'.join(lines))
'''
        result = subprocess.run(
            ['python3', '-c', code],
            capture_output=True, text=True, timeout=30
        )
        if result.returncode != 0:
            return {'status': 'error', 'error': result.stderr[:500]}
        text = Path(output_path).read_text(encoding='utf-8', errors='replace')
        return {'status': 'ok', 'chars': len(text)}
    except Exception as e:
        return {'status': 'error', 'error': str(e)}

# 格式 → 提取函数映射
EXTRACTORS = {
    '.pdf':  extract_pdf,
    '.pptx': extract_pptx,
    '.ppt':  extract_pptx,  # .ppt 也尝试 python-pptx
    '.docx': extract_docx,
    '.doc':  extract_doc,
    '.xlsx': extract_xlsx,
    '.xls':  extract_xlsx,
}

def extract_one(filepath, output_dir):
    """提取单个文件"""
    path = Path(filepath)
    ext = path.suffix.lower()
    output_path = os.path.join(output_dir, path.name + '.txt')
    file_hash = sha256(filepath)

    info = {
        'file': str(path),
        'name': path.name,
        'ext': ext,
        'sha256': file_hash,
        'size': path.stat().st_size,
    }

    extractor = EXTRACTORS.get(ext)
    if not extractor:
        info['status'] = 'skipped'
        info['note'] = f'不支持的文件格式: {ext}'
        return info

    result = extractor(str(path), output_path)
    info.update(result)
    info['output'] = output_path if result.get('status','').startswith('ok') else None
    return info

def extract_all(file_list, output_dir):
    """批量提取"""
    os.makedirs(output_dir, exist_ok=True)
    results = []
    for f in file_list:
        r = extract_one(f, output_dir)
        results.append(r)
        status = r.get('status','unknown')
        icon = '✅' if status.startswith('ok') else ('⚠️' if status == 'skipped' else '❌')
        print(f"  {icon} {r['name']} → {status}")
    return results

if __name__ == '__main__':
    if len(sys.argv) < 3:
        print("Usage: python3 extract.py <output_dir> <file1> [file2 ...]")
        print("       python3 extract.py <output_dir> @<list_file>")
        sys.exit(1)

    output_dir = sys.argv[1]

    if sys.argv[2].startswith('@'):
        with open(sys.argv[2][1:]) as f:
            files = [l.strip() for l in f if l.strip()]
    else:
        files = sys.argv[2:]

    results = extract_all(files, output_dir)

    # 输出 JSON 结果
    print('\n##RESULT_JSON##')
    print(json.dumps(results, ensure_ascii=False, indent=2))
